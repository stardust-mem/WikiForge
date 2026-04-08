"""PPTX 文档提取 — python-pptx"""

import io
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


@dataclass
class SlideData:
    slide_num: int
    title: str
    text_content: str
    notes: str
    images: list[bytes] = field(default_factory=list)
    image_types: list[str] = field(default_factory=list)


@dataclass
class PPTXExtractResult:
    slides: list[SlideData]
    total_slides: int
    total_text: str
    has_images: bool


def extract_pptx(file_path: Path) -> PPTXExtractResult:
    """
    提取 PPTX 内容：逐页提取标题 + 文字框 + 备注 + 图片

    输出格式化为 markdown：
    ## Slide 1: [标题]
    [文字内容]
    [备注：演讲者稿件]
    """
    prs = Presentation(str(file_path))
    slides = []
    all_text_parts = []

    for i, slide in enumerate(prs.slides):
        title = ""
        text_parts = []
        images = []
        image_types = []

        for shape in slide.shapes:
            # 标题
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    title = shape.text_frame.text.strip()
                else:
                    text = shape.text_frame.text.strip()
                    if text:
                        text_parts.append(text)

            # 表格
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    header = rows[0]
                    sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                    text_parts.append("\n".join([header, sep] + rows[1:]))

            # 图片
            if hasattr(shape, "image"):
                try:
                    img_bytes = shape.image.blob
                    if len(img_bytes) > 5000:
                        content_type = shape.image.content_type or "image/png"
                        images.append(img_bytes)
                        image_types.append(content_type)
                except Exception:
                    continue

        # 备注
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        text_content = "\n".join(text_parts)

        # 构建 slide markdown
        slide_md = f"## Slide {i+1}: {title}\n\n{text_content}"
        if notes:
            slide_md += f"\n\n> 备注：{notes}"
        all_text_parts.append(slide_md)

        slides.append(SlideData(
            slide_num=i + 1,
            title=title,
            text_content=text_content,
            notes=notes,
            images=images,
            image_types=image_types,
        ))

    total_text = "\n\n---\n\n".join(all_text_parts)
    has_images = any(len(s.images) > 0 for s in slides)

    return PPTXExtractResult(
        slides=slides,
        total_slides=len(slides),
        total_text=total_text,
        has_images=has_images,
    )
